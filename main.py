from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import shutil
import os
import tempfile
import base64
import time
import requests
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
import io

# Document processing libraries
from pptx import Presentation
from docx import Document
import pandas as pd
import fitz  # PyMuPDF for PDF processing

# PDF conversion libraries (currently not used - using text extraction instead)
# import pypandoc

# Additional document processing libraries
import csv
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, Alignment

# ReportLab for Excel to PDF conversion
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

# Matplotlib for Excel to PDF conversion
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend

app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with your frontend URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def create_text_image(text_content: str, page_number: int = 1, title: str = "Document") -> str:
    """Create a visual representation of text content as an image"""
    # Create a blank image
    width, height = 1200, 800
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)
    
    # Add title and page number
    draw.text((50, 30), f"{title} - Page {page_number}", fill='black')
    
    # Add text content with word wrapping
    y_position = 80
    words = text_content.split()
    current_line = ""
    
    for word in words:
        test_line = current_line + " " + word if current_line else word
        # Simple word wrapping (you can enhance this with proper text measurement)
        if len(test_line) > 80:  # Approximate characters per line
            if current_line:
                draw.text((50, y_position), current_line, fill='black')
                y_position += 20
                current_line = word
            else:
                draw.text((50, y_position), word, fill='black')
                y_position += 20
        else:
            current_line = test_line
    
    # Draw remaining text
    if current_line:
        draw.text((50, y_position), current_line, fill='black')
    
    # Convert to base64
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    img_base64 = base64.b64encode(buffer.getvalue()).decode()
    return img_base64

def process_word_document(file_path: str) -> list[str]:
    """Convert Word document to images via enhanced text extraction"""
    try:
        # Use enhanced text extraction for all Word formats
        print("  Processing Word document with enhanced text extraction")
        return process_word_document_enhanced(file_path)
        
    except Exception as e:
        print(f"  ERROR in Word processing: {str(e)}")
        print(f"  Error details: {type(e).__name__}: {str(e)}")
        # Fallback to basic text extraction if enhanced method fails
        print("  Falling back to basic text extraction method...")
        return process_word_document_fallback(file_path)

def process_word_document_enhanced(file_path: str) -> list[str]:
    """Enhanced Word document processing with better formatting"""
    doc = Document(file_path)
    images = []
    
    # Extract text from paragraphs with better formatting
    text_content = f"Word Document: {os.path.basename(file_path)}\n"
    text_content += "=" * 60 + "\n\n"
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            # Add paragraph formatting context
            if paragraph.style.name != 'Normal':
                text_content += f"[{paragraph.style.name}]: "
            text_content += paragraph.text + "\n\n"
    
    # Split content into pages (enhanced approach)
    words_per_page = 800  # More words per page for better readability
    words = text_content.split()
    pages = []
    
    for i in range(0, len(words), words_per_page):
        page_words = words[i:i + words_per_page]
        pages.append(" ".join(page_words))
    
    # Create images for each page
    for i, page_content in enumerate(pages):
        if page_content.strip():
            image = create_text_image(page_content, i + 1, "Word Document")
            images.append(image)
    
    return images

def process_word_document_fallback(file_path: str) -> list[str]:
    """Fallback method for Word document processing (old approach)"""
    doc = Document(file_path)
    images = []
    
    # Extract text from paragraphs
    text_content = ""
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_content += paragraph.text + "\n\n"
    
    # Split content into pages (simple approach)
    words_per_page = 500
    words = text_content.split()
    pages = []
    
    for i in range(0, len(words), words_per_page):
        page_words = words[i:i + words_per_page]
        pages.append(" ".join(page_words))
    
    # Create images for each page
    for i, page_content in enumerate(pages):
        if page_content.strip():
            image = create_text_image(page_content, i + 1, "Word Document")
            images.append(image)
    
    return images

def process_excel_document(file_path: str) -> list[str]:
    """Convert Excel/CSV document to images via PDF conversion using matplotlib"""
    try:
        # Convert Excel/CSV to PDF first using matplotlib approach
        pdf_path = file_path.replace('.xlsx', '_converted.pdf').replace('.xls', '_converted.pdf').replace('.csv', '_converted.pdf').replace('.ods', '_converted.pdf')
        print(f"  Converting Excel/CSV document to PDF: {file_path} -> {pdf_path}")
        
        # Convert to PDF using matplotlib + reportlab
        convert_excel_to_pdf_matplotlib(file_path, pdf_path)
        
        # Process the PDF using the existing PDF processing function
        print(f"  Processing converted PDF: {pdf_path}")
        images = process_pdf_document(pdf_path)
        
        # Clean up the temporary PDF file
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
            print(f"  Cleaned up temporary PDF: {pdf_path}")
        
        return images
        
    except Exception as e:
        print(f"  ERROR in Excel/CSV to PDF conversion: {str(e)}")
        print(f"  Error details: {type(e).__name__}: {str(e)}")
        # Fallback to enhanced text extraction if PDF conversion fails
        print("  Falling back to enhanced text extraction method...")
        return process_excel_document_enhanced(file_path)

def df_to_image(df, out_png_path, page_width_in, page_height_in, dpi=200, title=None):
    """
    Render a DataFrame to a PNG image sized to fit within the given
    width/height (in inches). We compute a font size that fits rows/cols.
    """
    # margins inside the image (in inches)
    left_margin_in = 0.25
    right_margin_in = 0.25
    top_margin_in = 0.5 if title else 0.25
    bottom_margin_in = 0.25

    usable_w_in = max(1e-3, page_width_in - left_margin_in - right_margin_in)
    usable_h_in = max(1e-3, page_height_in - top_margin_in - bottom_margin_in)

    # Create the figure sized to the full page minus margins
    fig_w_in = page_width_in
    fig_h_in = page_height_in

    fig, ax = plt.subplots(figsize=(fig_w_in, fig_h_in), dpi=dpi)
    ax.axis("off")

    # Positioning rect for the table (in axes fraction). We'll leave margins.
    # Convert inches to axis fraction.
    left = left_margin_in / fig_w_in
    right = 1.0 - (right_margin_in / fig_w_in)
    bottom = bottom_margin_in / fig_h_in
    top = 1.0 - (top_margin_in / fig_h_in)

    # Calculate a font size that (roughly) fits rows/columns into usable area.
    # Text size in points; 72 points = 1 inch.
    nrows, ncols = df.shape
    nrows = max(1, nrows)
    ncols = max(1, ncols)

    # Heuristic: each row needs about 1.3 * (fontsize/72) inches in height;
    # each column needs about 2.3 * (fontsize/72) inches in width.
    # Tune factors to balance readability vs. fitting large sheets.
    if nrows == 0:
        # Empty sheet—make a tiny table with just headers (if any)
        nrows_eff = 1
    else:
        nrows_eff = nrows + 1  # account for header row

    # Height-limited font size
    fs_by_h = (usable_h_in * 72) / (1.3 * nrows_eff)
    # Width-limited font size
    fs_by_w = (usable_w_in * 72) / (2.3 * max(1, ncols))

    fontsize = max(5, min(12, fs_by_h, fs_by_w))

    # Prepare table data with headers
    data = [list(df.columns)]
    if nrows > 0:
        # Convert all cells to string to avoid weird rendering
        body = df.astype(str).values.tolist()
        data.extend(body)

    # Draw title (if any)
    if title:
        ax.text(
            0.5, 1.0 - (top_margin_in - 0.25) / fig_h_in,  # position slightly above table area
            title,
            ha="center",
            va="top",
            fontsize=min(16, max(10, fontsize + 2)),
            fontweight="bold",
            transform=ax.transAxes,
        )

    # Create the table
    table = ax.table(
        cellText=data,
        loc="center",
        colLabels=None,  # we already placed header as first row of cellText
        cellLoc="center",
        bbox=(left, bottom, right - left, top - bottom)  # (x, y, w, h) in axes fraction
    )

    # Style header row
    header_row_idx = 0
    for (row, col), cell in table.get_celld().items():
        cell.set_linewidth(0.3)
        if row == header_row_idx:
            cell.set_text_props(fontweight="bold")
    # Set all font sizes
    for key, cell in table.get_celld().items():
        cell.set_fontsize(fontsize)

    # Tight layout off; we explicitly controlled bbox.
    fig.savefig(out_png_path, bbox_inches="tight", dpi=dpi)
    plt.close(fig)

def convert_excel_to_pdf_matplotlib(file_path: str, pdf_path: str):
    """Convert Excel/CSV file to PDF using matplotlib + reportlab (one page per sheet)"""
    try:
        # Page settings
        page_w_pt, page_h_pt = A4
        margin_pt = 36  # 0.5 inch
        content_w_pt = page_w_pt - 2 * margin_pt
        content_h_pt = page_h_pt - 2 * margin_pt

        # Convert points to inches for our image renderer
        content_w_in = content_w_pt / 72.0
        content_h_in = content_h_pt / 72.0

        c = canvas.Canvas(pdf_path, pagesize=A4)

        if file_path.endswith('.csv'):
            # Handle CSV files
            df = pd.read_csv(file_path)
            with tempfile.TemporaryDirectory() as td:
                img_path = os.path.join(td, "sheet.png")
                # Render DF to PNG sized to fit the content area
                title = f"{os.path.basename(file_path)}"
                df_to_image(df, img_path, content_w_in, content_h_in, dpi=200, title=title)

                # Draw on PDF page (centered in content area)
                from PIL import Image
                img = Image.open(img_path)
                img_w_px, img_h_px = img.size
                img_dpi = 200.0  # must match df_to_image dpi
                img_w_in = img_w_px / img_dpi
                img_h_in = img_h_px / img_dpi
                img_w_pt = img_w_in * 72.0
                img_h_pt = img_h_in * 72.0

                scale = min(content_w_pt / img_w_pt, content_h_pt / img_h_pt)
                draw_w = img_w_pt * scale
                draw_h = img_h_pt * scale

                x = (page_w_pt - draw_w) / 2.0
                y = (page_h_pt - draw_h) / 2.0

                c.drawImage(img_path, x, y, width=draw_w, height=draw_h, preserveAspectRatio=True, anchor='c')

            c.showPage()
        else:
            # Handle Excel files
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = xls.parse(sheet_name=sheet_name, dtype=str)

                with tempfile.TemporaryDirectory() as td:
                    img_path = os.path.join(td, "sheet.png")
                    # Render DF to PNG sized to fit the content area
                    title = f"{os.path.basename(file_path)} — {sheet_name}"
                    df_to_image(df, img_path, content_w_in, content_h_in, dpi=200, title=sheet_name)

                    # Draw on PDF page (centered in content area)
                    from PIL import Image
                    img = Image.open(img_path)
                    img_w_px, img_h_px = img.size
                    img_dpi = 200.0  # must match df_to_image dpi
                    img_w_in = img_w_px / img_dpi
                    img_h_in = img_h_px / img_dpi
                    img_w_pt = img_w_in * 72.0
                    img_h_pt = img_h_in * 72.0

                    scale = min(content_w_pt / img_w_pt, content_h_pt / img_h_pt)
                    draw_w = img_w_pt * scale
                    draw_h = img_h_pt * scale

                    x = (page_w_pt - draw_w) / 2.0
                    y = (page_h_pt - draw_h) / 2.0

                    c.drawImage(img_path, x, y, width=draw_w, height=draw_h, preserveAspectRatio=True, anchor='c')

                c.showPage()

        c.save()
        print(f"  Successfully converted Excel/CSV to PDF using matplotlib: {pdf_path}")
        return True
        
    except Exception as e:
        print(f"  ERROR in Excel to PDF conversion (matplotlib): {str(e)}")
        raise e

def convert_excel_to_pdf(file_path: str, pdf_path: str):
    """Convert Excel/CSV file to PDF using openpyxl + reportlab"""
    try:
        # Create PDF document
        doc = SimpleDocTemplate(pdf_path, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
        story = []
        styles = getSampleStyleSheet()
        
        if file_path.endswith('.csv'):
            # Handle CSV files
            df = pd.read_csv(file_path)
            
            # Add file title
            title = Paragraph(f"CSV File: {os.path.basename(file_path)}", styles['Heading1'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Add file info
            info_text = f"Rows: {len(df)}, Columns: {len(df.columns)}"
            info = Paragraph(info_text, styles['Normal'])
            story.append(info)
            story.append(Spacer(1, 12))
            
            # Convert DataFrame to table
            data = [df.columns.tolist()] + df.values.tolist()
            
            # Limit rows for PDF (first 50 rows)
            if len(data) > 51:  # 1 header + 50 data rows
                data = data[:51]
                data.append(["...", "...", "...", "...", "..."] * (len(df.columns) // 5 + 1))
            
            # Create table
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ]))
            story.append(table)
            
        else:
            # Handle Excel files
            workbook = load_workbook(file_path, data_only=True)
            
            for sheet_idx, sheet_name in enumerate(workbook.sheetnames):
                if sheet_idx > 0:
                    story.append(PageBreak())
                
                worksheet = workbook[sheet_name]
                
                # Add sheet title
                title = Paragraph(f"Sheet: {sheet_name}", styles['Heading1'])
                story.append(title)
                story.append(Spacer(1, 12))
                
                # Get data from worksheet
                data = []
                max_cols = 0
                
                for row in worksheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    data.append(row_data)
                    max_cols = max(max_cols, len(row_data))
                
                # Pad rows to have same number of columns
                for row in data:
                    while len(row) < max_cols:
                        row.append("")
                
                # Limit rows for PDF (first 30 rows per sheet)
                if len(data) > 31:  # 1 header + 30 data rows
                    data = data[:31]
                    data.append(["..."] * max_cols)
                
                # Create table
                if data:
                    table = Table(data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        print(f"  Successfully converted Excel/CSV to PDF: {pdf_path}")
        return True
        
    except Exception as e:
        print(f"  ERROR in Excel to PDF conversion: {str(e)}")
        raise e

def process_excel_document_enhanced(file_path: str) -> list[str]:
    """Enhanced Excel/CSV document processing with better formatting (fallback)"""
    images = []
    
    if file_path.endswith('.csv'):
        # Handle CSV files with better formatting
        df = pd.read_csv(file_path)
        text_content = f"CSV File: {os.path.basename(file_path)}\n"
        text_content += "=" * 60 + "\n\n"
        text_content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n\n"
        text_content += "Column Headers:\n"
        text_content += "- " + "\n- ".join(df.columns.tolist()) + "\n\n"
        text_content += "Data Preview:\n"
        text_content += df.head(20).to_string(index=False, max_cols=10)
        
        image = create_text_image(text_content, 1, "CSV File")
        images.append(image)
    else:
        # Handle Excel files with better formatting
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_idx, sheet_name in enumerate(excel_file.sheet_names):
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            text_content = f"Excel Sheet: {sheet_name}\n"
            text_content += "=" * 60 + "\n\n"
            text_content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n\n"
            text_content += "Column Headers:\n"
            text_content += "- " + "\n- ".join(df.columns.tolist()) + "\n\n"
            text_content += "Data Preview:\n"
            text_content += df.head(15).to_string(index=False, max_cols=8)
            
            image = create_text_image(text_content, sheet_idx + 1, f"Excel - {sheet_name}")
            images.append(image)
    
    return images

def process_excel_document_fallback(file_path: str) -> list[str]:
    """Fallback method for Excel/CSV document processing (old approach)"""
    images = []
    
    if file_path.endswith('.csv'):
        # Handle CSV files
        df = pd.read_csv(file_path)
        text_content = f"CSV File\n\n"
        text_content += df.to_string(index=False)
        image = create_text_image(text_content, 1, "CSV File")
        images.append(image)
    else:
        # Handle Excel files
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert DataFrame to text representation
            text_content = f"Sheet: {sheet_name}\n\n"
            text_content += df.to_string(index=False)
            
            image = create_text_image(text_content, 1, f"Excel - {sheet_name}")
            images.append(image)
    
    return images

def process_powerpoint_document(file_path: str) -> list[str]:
    """Convert PowerPoint document to images via enhanced text extraction"""
    try:
        # Use enhanced text extraction since pypandoc doesn't support PowerPoint
        print("  Processing PowerPoint document with enhanced text extraction (pypandoc doesn't support .pptx)")
        return process_powerpoint_document_enhanced(file_path)
        
    except Exception as e:
        print(f"  ERROR in PowerPoint processing: {str(e)}")
        print(f"  Error details: {type(e).__name__}: {str(e)}")
        # Fallback to basic text extraction if enhanced method fails
        print("  Falling back to basic text extraction method...")
        return process_powerpoint_document_fallback(file_path)

def process_powerpoint_document_enhanced(file_path: str) -> list[str]:
    """Enhanced PowerPoint document processing with better formatting"""
    prs = Presentation(file_path)
    images = []
    
    for i, slide in enumerate(prs.slides):
        # Extract text from shapes with better formatting
        text_content = f"Slide {i + 1}\n"
        text_content += "=" * 50 + "\n\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Add shape type context
                shape_type = str(shape.shape_type).split('.')[-1] if hasattr(shape, 'shape_type') else "Text"
                text_content += f"[{shape_type}]: {shape.text.strip()}\n\n"
        
        if text_content.strip():
            image = create_text_image(text_content, i + 1, "PowerPoint")
            images.append(image)
    
    return images

def process_powerpoint_document_fallback(file_path: str) -> list[str]:
    """Fallback method for PowerPoint document processing (old approach)"""
    prs = Presentation(file_path)
    images = []
    
    for i, slide in enumerate(prs.slides):
        # Extract text from shapes
        text_content = f"Slide {i + 1}\n\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content += shape.text + "\n"
        
        image = create_text_image(text_content, i + 1, "PowerPoint")
        images.append(image)
    
    return images

def process_pdf_document(file_path: str) -> list[str]:
    """Convert PDF document to images"""
    # Open PDF with PyMuPDF
    pdf_document = fitz.open(file_path)
    images_base64 = []
    
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        
        # Convert page to image with higher resolution
        mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to PIL Image
        img_data = pix.tobytes("png")
        img = Image.open(BytesIO(img_data))
        
        # Convert to base64
        buffer = BytesIO()
        img.save(buffer, format='PNG')
        img_base64 = base64.b64encode(buffer.getvalue()).decode()
        images_base64.append(img_base64)
    
    pdf_document.close()
    return images_base64

def process_text_document(file_path: str) -> list[str]:
    """Convert text document to images"""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Split content into pages
    words_per_page = 500
    words = content.split()
    pages = []
    
    for i in range(0, len(words), words_per_page):
        page_words = words[i:i + words_per_page]
        pages.append(" ".join(page_words))
    
    images = []
    for i, page_content in enumerate(pages):
        if page_content.strip():
            image = create_text_image(page_content, i + 1, "Text Document")
            images.append(image)
    
    return images

@app.post("/document-to-images")
async def document_to_images(file: UploadFile = File(...)):
    # Add detailed logging
    print(f"Received file upload request:")
    print(f"  Filename: {file.filename}")
    print(f"  Content type: {file.content_type}")
    print(f"  File size: {file.size if hasattr(file, 'size') else 'unknown'}")
    
    # Validate file type
    filename = file.filename.lower() if file.filename else ""
    supported_extensions = ['.pdf', '.docx', '.doc', '.rtf', '.odt', '.xlsx', '.xls', '.ods', '.csv', '.pptx', '.ppt', '.odp', '.txt']
    
    print(f"  File extension check: {filename}")
    print(f"  Supported extensions: {supported_extensions}")
    
    if not filename:
        print("  ERROR: No filename provided")
        raise HTTPException(
            status_code=400, 
            detail="No filename provided"
        )
    
    if not any(filename.endswith(ext) for ext in supported_extensions):
        print(f"  ERROR: Unsupported file type: {filename}")
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: {filename}. Supported: {', '.join(supported_extensions)}"
        )
    
    print(f"  File validation passed, processing...")
    
    # Save uploaded file to a temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
        shutil.copyfileobj(file.file, tmp)
        tmp_path = tmp.name
    
    try:
        # Process based on file type
        if filename.endswith('.pdf'):
            print("  Processing as PDF...")
            images = process_pdf_document(tmp_path)
        elif filename.endswith(('.docx', '.doc', '.rtf', '.odt')):
            print("  Processing as Word document...")
            images = process_word_document(tmp_path)
        elif filename.endswith(('.xlsx', '.xls', '.ods', '.csv')):
            print("  Processing as Excel/CSV document...")
            images = process_excel_document(tmp_path)
        elif filename.endswith(('.pptx', '.ppt', '.odp')):
            print("  Processing as PowerPoint document...")
            images = process_powerpoint_document(tmp_path)
        elif filename.endswith('.txt'):
            print("  Processing as text document...")
            images = process_text_document(tmp_path)
        else:
            print(f"  ERROR: Unsupported file type after validation: {filename}")
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        print(f"  Processing completed successfully. Generated {len(images)} images.")
        
        return {
            "images": images, 
            "num_pages": len(images),
            "file_type": os.path.splitext(filename)[1][1:].upper()
        }
        
    except Exception as e:
        print(f"  ERROR during processing: {str(e)}")
        print(f"  Error type: {type(e).__name__}")
        print(f"  Error details: {str(e)}")
        import traceback
        print(f"  Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

@app.post("/images-to-images")
async def images_to_images(files: list[UploadFile] = File(...)):
    """Handle multiple image uploads as a document sequence"""
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    
    # Validate that all files are images
    supported_image_types = ['image/jpeg', 'image/jpg', 'image/png', 'image/gif', 'image/webp', 'image/bmp']
    
    for file in files:
        if file.content_type not in supported_image_types:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported image type: {file.content_type}. Supported: {', '.join(supported_image_types)}"
            )
    
    images_base64 = []
    
    try:
        for i, file in enumerate(files):
            # Read image file
            image_bytes = await file.read()
            
            # Convert to PIL Image for processing
            img = Image.open(BytesIO(image_bytes))
            
            # Convert to base64
            buffer = BytesIO()
            img.save(buffer, format='PNG')
            img_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            images_base64.append(img_base64)
        
        return {
            "images": images_base64,
            "num_pages": len(images_base64),
            "file_type": "IMAGE_SEQUENCE",
            "message": f"Processed {len(images_base64)} images as document sequence"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing images: {str(e)}")

@app.post("/convert-images-to-base64")
async def convert_images_to_base64(
    images: list[str],  # Base64 strings from the request
    processing_mode: str = "full",
    selected_pages: list[int] = None
):
    """
    Convert base64 image strings and handle page selection
    """
    try:
        print(f"Received image conversion request:")
        print(f"  Number of images: {len(images)}")
        print(f"  Processing mode: {processing_mode}")
        print(f"  Selected pages: {selected_pages}")
        
        if not images:
            raise HTTPException(status_code=400, detail="No images provided")
        
        # Validate that all images are valid base64 strings
        images_base64 = []
        for i, image_base64 in enumerate(images):
            try:
                # Validate base64 string
                if not image_base64 or not isinstance(image_base64, str):
                    raise ValueError("Invalid base64 string")
                
                # Test if it's valid base64 by trying to decode
                import base64
                base64.b64decode(image_base64)
                
                images_base64.append(image_base64)
                print(f"  Validated image {i+1}/{len(images)} as base64")
            except Exception as e:
                print(f"  Error validating image {i+1}: {str(e)}")
                raise HTTPException(status_code=400, detail=f"Invalid base64 image {i+1}: {str(e)}")
        
        print(f"  Successfully validated {len(images_base64)} images as base64")
        
        # Handle page selection
        if processing_mode == "selection" and selected_pages:
            # Validate selected pages
            if not all(1 <= p <= len(images_base64) for p in selected_pages):
                raise HTTPException(
                    status_code=400, 
                    detail=f"Invalid page selection. Valid range: 1-{len(images_base64)}"
                )
            
            # Filter to selected pages (1-indexed to 0-indexed)
            selected_images = []
            for page_num in selected_pages:
                image_index = page_num - 1  # Convert 1-indexed to 0-indexed
                if 0 <= image_index < len(images_base64):
                    selected_images.append(images_base64[image_index])
            
            images_base64 = selected_images
            print(f"  Filtered to {len(images_base64)} selected pages: {selected_pages}")
        
        return {
            "images_base64": images_base64,
            "num_pages": len(images_base64),
            "selected_pages": selected_pages if processing_mode == "selection" else None,
            "processing_mode": processing_mode,
            "message": f"Successfully processed {len(images_base64)} images"
        }
        
    except Exception as e:
        print(f"  ERROR during image conversion: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting images: {str(e)}")

@app.get("/")
def root():
    return {
        "message": "Unified Document to Image Microservice - Integrated with document-base64-analyzer",
        "supported_formats": ["PDF", "DOCX", "DOC", "XLSX", "XLS", "PPTX", "PPT", "TXT"],
        "endpoints": {
            "document_to_images": "/document-to-images",
            "images_to_images": "/images-to-images",
            "process_with_ai": "/process-with-ai",
            "convert_images_to_base64": "/convert-images-to-base64"
        },
        "features": [
            "Single document upload (PDF, Word, Excel, PowerPoint, Text)",
            "Multiple image uploads as document sequence",
            "Returns base64 images ready for AI analysis",
            "Direct AI processing with page selection support",
            "Integrated with document-base64-analyzer.onrender.com",
            "Efficient page selection - only selected pages sent to AI",
            "Image conversion from downloaded data to base64",
            "Page selection filtering for partial document processing"
        ]
    }

@app.get("/health")
def health_check():
    return {
        "status": "healthy",
        "service": "Unified Document to Image Microservice - Integrated with document-base64-analyzer",
        "timestamp": "2024-01-01T00:00:00Z",
        "ai_microservice": "document-base64-analyzer.onrender.com",
        "integration_status": "active"
    }

@app.post("/process-with-ai")
async def process_with_ai(
    file: UploadFile = File(...),
    processing_mode: str = "full",  # "full", "smart", or "selection"
    selected_pages: list[int] = None
):
    """
    Convert document to images and process with AI microservice based on user choice
    """
    print(f"Received AI processing request:")
    print(f"  Filename: {file.filename}")
    print(f"  Processing mode: {processing_mode}")
    print(f"  Selected pages: {selected_pages}")
    
    # First, convert document to images using existing logic
    try:
        # Save uploaded file to a temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as tmp:
            shutil.copyfileobj(file.file, tmp)
            tmp_path = tmp.name
        
        # Process based on file type (reuse existing logic)
        filename = file.filename.lower() if file.filename else ""
        if filename.endswith('.pdf'):
            images = process_pdf_document(tmp_path)
        elif filename.endswith(('.docx', '.doc')):
            images = process_word_document(tmp_path)
        elif filename.endswith(('.xlsx', '.xls')):
            images = process_excel_document(tmp_path)
        elif filename.endswith(('.pptx', '.ppt')):
            images = process_powerpoint_document(tmp_path)
        elif filename.endswith('.txt'):
            images = process_text_document(tmp_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        os.remove(tmp_path)
        
        print(f"  Document converted to {len(images)} images successfully")
        
        # Now call the document-base64-analyzer microservice
        ai_microservice_url = "https://document-base64-analyzer.onrender.com"
        
        # Prepare the payload for AI microservice
        payload = {
            "job_id": f"job_{int(time.time())}",
            "user_id": "user_123",  # This should come from the request
            "images_base64": images,
            "num_pages": len(images),
            "file_type": os.path.splitext(filename)[1][1:].upper(),
            "fallback_text": f"Document: {file.filename}"
        }
        
        # Handle page selection by filtering images
        if processing_mode == "selection" and selected_pages:
            # Validate selected pages
            if not all(1 <= p <= len(images) for p in selected_pages):
                raise HTTPException(
                    status_code=400, 
                    detail=f"Invalid page selection. Valid range: 1-{len(images)}"
                )
            
            # Filter images to only selected pages (1-indexed to 0-indexed)
            selected_images = []
            for page_num in selected_pages:
                image_index = page_num - 1  # Convert 1-indexed to 0-indexed
                if 0 <= image_index < len(images):
                    selected_images.append(images[image_index])
            
            # Update payload with only selected images
            payload["images_base64"] = selected_images
            payload["num_pages"] = len(selected_images)
            payload["selected_pages"] = selected_pages
            
            print(f"  Processing {len(selected_images)} selected pages: {selected_pages}")
        else:
            # Full processing - use all images
            payload["num_pages"] = len(images)
            print(f"  Processing all {len(images)} pages")
        
        # Use single endpoint for all processing modes
        ai_endpoint = f"{ai_microservice_url}/process-document"
        
        print(f"  Calling AI microservice: {ai_endpoint}")
        
        # Make the request to AI microservice
        import requests
        response = requests.post(ai_endpoint, json=payload, timeout=30)
        
        if response.status_code == 200:
            ai_result = response.json()
            print(f"  AI microservice response: {ai_result}")
            
            return {
                "status": "success",
                "message": "Document sent to AI processing successfully",
                "ai_response": ai_result,
                "document_info": {
                    "images_generated": len(images),
                    "file_type": os.path.splitext(filename)[1][1:].upper(),
                    "processing_mode": processing_mode
                }
            }
        else:
            print(f"  AI microservice error: {response.status_code} - {response.text}")
            raise HTTPException(
                status_code=500, 
                detail=f"AI microservice error: {response.status_code}"
            )
            
    except Exception as e:
        print(f"  ERROR during AI processing: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}") 